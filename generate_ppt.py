"""
Generate DPS Presentation with dark theme and smooth transitions.
"""
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# --- Constants ---
BG_COLOR = RGBColor(0x0A, 0x0E, 0x17)
ACCENT_BLUE = RGBColor(0x4D, 0x8B, 0xFF)
ACCENT_CYAN = RGBColor(0x22, 0xD3, 0xEE)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
TEXT_PRIMARY = RGBColor(0xE8, 0xEC, 0xF4)
TEXT_SECONDARY = RGBColor(0x88, 0x99, 0xB8)
TEXT_MUTED = RGBColor(0x55, 0x66, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0x0E, 0x14, 0x21)

with open("models/model_meta.json") as f:
    meta = json.load(f)
metrics = meta["metrics"]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_transition(slide):
    sld_elem = slide._element
    trans = sld_elem.makeelement(qn('p:transition'), {
        'spd': 'slow', 'advClick': '1'
    })
    fade = trans.makeelement(qn('p:fade'), {'thruBlk': '1'})
    trans.append(fade)
    sld_elem.append(trans)

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_heading(slide, text):
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                text, font_size=32, color=ACCENT_BLUE, bold=True)
    # Underline bar
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.05), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_CYAN
    line.line.fill.background()

def add_footer(slide, is_first_page=False):
    if is_first_page:
        add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.5),
                    "Kathan  •  Aniket  •  Rudra", font_size=12, color=TEXT_MUTED,
                    alignment=PP_ALIGN.CENTER)

def add_bullet_slide(slide, items, start_top=1.5, left=0.8, font_size=18):
    for i, item in enumerate(items):
        add_textbox(slide, Inches(left), Inches(start_top + i * 0.55),
                    Inches(10), Inches(0.5), f"•  {item}", font_size=font_size,
                    color=TEXT_PRIMARY)

# ===================== SLIDE 1: TITLE =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_transition(slide)

# Accent bar at top
bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
            "Diabetes Prediction System [DPS]", font_size=48, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
            "An End-to-End Machine Learning Web Application for Diabetes Risk Assessment",
            font_size=20, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Decorative line
deco = slide.shapes.add_shape(1, Inches(5), Inches(4.3), Inches(3.333), Pt(2))
deco.fill.solid()
deco.fill.fore_color.rgb = ACCENT_CYAN
deco.line.fill.background()

add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
            "Mini Project Presentation", font_size=16, color=TEXT_MUTED,
            alignment=PP_ALIGN.CENTER)
add_footer(slide, is_first_page=True)

# ===================== SLIDE 2: INTRODUCTION =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_transition(slide)
add_heading(slide, "Introduction")
items = [
    "DPS is a lightweight, end-to-end Machine Learning project for predicting diabetes risk.",
    "Built using Python, Scikit-learn, and deployed as a static web application.",
    "Uses the Pima Indians Diabetes Dataset (768 records, 8 clinical features).",
    "Trains 4 traditional ML models and selects the best performer automatically.",
    "Features real-time browser-based inference — no server required after export.",
    "Includes explainable AI with feature contribution analysis for every prediction.",
    "Designed for educational and portfolio demonstration purposes."
]
add_bullet_slide(slide, items)
add_footer(slide)

# ===================== SLIDE 3: TECH STACK =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_transition(slide)
add_heading(slide, "Technology Stack")

categories = [
    ("ML & Data Science", "Python, Scikit-learn, Pandas, NumPy, Imbalanced-learn (SMOTE)"),
    ("Model Export", "m2cgen — converts trained models to pure JavaScript"),
    ("Frontend", "HTML5, CSS3, Vanilla JavaScript, Chart.js"),
    ("Preprocessing", "StandardScaler, Mutual Information, Median Imputation"),
    ("Deployment", "GitHub Pages (static site), no backend needed"),
    ("Version Control", "Git & GitHub"),
]
for i, (cat, desc) in enumerate(categories):
    y = 1.5 + i * 0.75
    add_textbox(slide, Inches(0.8), Inches(y), Inches(3.5), Inches(0.5),
                cat, font_size=18, color=ACCENT_CYAN, bold=True)
    add_textbox(slide, Inches(4.5), Inches(y), Inches(8), Inches(0.5),
                desc, font_size=16, color=TEXT_PRIMARY)
add_footer(slide)

# ===================== SLIDE 4: DATASET PREVIEW =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_transition(slide)
add_heading(slide, "Dataset Preview — First 5 Rows")
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.5),
            "The Pima Indians Diabetes Dataset contains 768 samples with 8 clinical features and 1 binary target (Outcome).",
            font_size=14, color=TEXT_SECONDARY)
try:
    slide.shapes.add_picture("dataset_preview.png", Inches(0.5), Inches(2.0), Inches(12.3))
except:
    add_textbox(slide, Inches(1), Inches(3), Inches(10), Inches(1),
                "[Dataset preview image not found]", font_size=16, color=TEXT_MUTED)
add_footer(slide)

# ===================== SLIDE 5: PREPROCESSING =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_transition(slide)
add_heading(slide, "Data Preprocessing Pipeline")
steps = [
    "Step 1: Data Loading — Read the CSV file with Pandas (768 rows × 9 columns).",
    "Step 2: Missing Value Imputation — Replace 0s in Glucose, BP, SkinThickness, Insulin, BMI with column medians.",
    "Step 3: Feature Importance Analysis — Rank features using Mutual Information scores.",
    "Step 4: Train-Test Split — 80:20 stratified split preserving class distribution.",
    "Step 5: Feature Scaling — StandardScaler normalizes all features to zero mean and unit variance.",
    "Step 6: Class Balancing — SMOTE generates synthetic minority samples to fix class imbalance.",
    "Step 7: Output — Returns balanced training data, scaled test data, scaler, and feature names."
]
add_bullet_slide(slide, steps, start_top=1.4, font_size=16)
add_footer(slide)

# ===================== SLIDE 6: FEATURE IMPORTANCE =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_transition(slide)
add_heading(slide, "Feature Importance Analysis")
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.5),
            "Mutual Information scores reveal which clinical features have the strongest predictive power for diabetes.",
            font_size=14, color=TEXT_SECONDARY)
try:
    slide.shapes.add_picture("feature_importance.png", Inches(1.5), Inches(1.9), Inches(10))
except:
    pass
add_footer(slide)

# ===================== SLIDE 7: HOW THE MODEL WORKS =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_transition(slide)
add_heading(slide, "How the Model Works")
steps = [
    "1. User enters 8 clinical values (Glucose, BMI, Age, etc.) into the web form.",
    "2. The JavaScript frontend applies StandardScaler normalization using saved mean/scale values.",
    "3. Scaled features are passed through the exported Random Forest decision tree logic (pure JS).",
    "4. The model outputs a probability score between 0% and 100% for diabetes risk.",
    "5. Risk is categorized: Low (<30%), Moderate (30–60%), or High (>60%).",
    "6. Feature contributions are computed using z-score deviation analysis.",
    "7. Results are displayed with an animated gauge, verdict badge, and explanation bars."
]
add_bullet_slide(slide, steps, start_top=1.4, font_size=16)
add_footer(slide)

# ===================== SLIDE 8: ML MODELS USED =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_transition(slide)
add_heading(slide, "ML Models & Accuracy Scores")

headers = ["Model", "Train Acc", "Test Acc", "Precision", "Recall", "F1-Score", "ROC-AUC"]
col_widths = [2.2, 1.5, 1.5, 1.5, 1.5, 1.5, 1.5]
table_width = sum(col_widths)
left = (13.333 - table_width) / 2

tbl = slide.shapes.add_table(len(metrics) + 1, len(headers),
                              Inches(left), Inches(1.8),
                              Inches(table_width), Inches(0.5 * (len(metrics) + 1)))
table = tbl.table

for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_BLUE

for r, m in enumerate(sorted(metrics, key=lambda x: x["Accuracy"], reverse=True)):
    vals = [m["Model"], f"{m.get('Train Accuracy','N/A')}%", f"{m['Accuracy']}%",
            f"{m['Precision']}%", f"{m['Recall']}%", f"{m['F1-Score']}%", f"{m.get('ROC-AUC','N/A')}%"]
    for c, v in enumerate(vals):
        cell = table.cell(r + 1, c)
        cell.text = str(v)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_PRIMARY
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x12, 0x18, 0x2A) if r % 2 == 0 else CARD_BG

best = max(metrics, key=lambda x: x["Accuracy"])
add_textbox(slide, Inches(left), Inches(4.6), Inches(table_width), Inches(0.5),
            f"★  Best Model: {best['Model']} — {best['Accuracy']}% Test Accuracy",
            font_size=20, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_footer(slide)

# ===================== SLIDE 9: TRAIN vs TEST ACCURACY GRAPH =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_transition(slide)
add_heading(slide, "Train vs Test Accuracy")
try:
    slide.shapes.add_picture("train_test_accuracy.png", Inches(1.5), Inches(1.5), Inches(10))
except:
    pass
add_footer(slide)

# ===================== SLIDE 10: ML MODELS ACCURACY GRAPH =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_transition(slide)
add_heading(slide, "ML Models Performance Metrics")
try:
    slide.shapes.add_picture("ml_models_accuracy.png", Inches(1.2), Inches(1.5), Inches(10.5))
except:
    pass
add_footer(slide)

# ===================== SLIDE 11: WEB INTERFACE =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_transition(slide)
add_heading(slide, "Web Application Interface")
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.5),
            "The live web app runs entirely in the browser using a JavaScript-exported ML model. No server required.",
            font_size=14, color=TEXT_SECONDARY)
try:
    slide.shapes.add_picture("website_ui.png", Inches(2.5), Inches(2.0), Inches(8))
except:
    add_textbox(slide, Inches(1), Inches(3), Inches(10), Inches(1),
                "[Website UI image not found]", font_size=16, color=TEXT_MUTED)
add_footer(slide)

# ===================== SLIDE 12: THANK YOU =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_transition(slide)
add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
            "Thank You!", font_size=54, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
deco = slide.shapes.add_shape(1, Inches(5), Inches(3.7), Inches(3.333), Pt(2))
deco.fill.solid()
deco.fill.fore_color.rgb = ACCENT_CYAN
deco.line.fill.background()
add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
            "Diabetes Prediction System [DPS]", font_size=20, color=ACCENT_BLUE,
            alignment=PP_ALIGN.CENTER)
add_footer(slide)

# Save
output_path = "DPS_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
